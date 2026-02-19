import io
import re
from typing import Optional
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
import config


def _chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> list[str]:
    """Split text into overlapping chunks."""
    chunk_size = chunk_size or config.MAX_CHUNK_SIZE
    overlap = overlap or config.CHUNK_OVERLAP
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks


def extract_youtube_id(url: str) -> Optional[str]:
    """Extract YouTube video ID from various URL formats."""
    patterns = [
        r"(?:v=|\/)([0-9A-Za-z_-]{11}).*",
        r"(?:youtu\.be\/)([0-9A-Za-z_-]{11})",
        r"(?:embed\/)([0-9A-Za-z_-]{11})",
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


class DocumentProcessor:
    """Handles extraction and chunking of text from all supported sources."""

    def process_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF bytes."""
        text = ""
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()

    def process_docx(self, file_bytes: bytes) -> str:
        """Extract text from DOCX bytes."""
        doc = DocxDocument(io.BytesIO(file_bytes))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs).strip()

    def process_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PPTX bytes."""
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts).strip()

    def process_txt(self, file_bytes: bytes) -> str:
        """Extract text from plain text / notepad files."""
        try:
            return file_bytes.decode("utf-8").strip()
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1").strip()

    def process_youtube(self, url: str) -> str:
        """Extract transcript from a YouTube video URL."""
        video_id = extract_youtube_id(url)
        if not video_id:
            raise ValueError("Invalid YouTube URL. Please check and try again.")
        try:
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            text = " ".join([entry["text"] for entry in transcript_list])
            return text.strip()
        except Exception as e:
            raise RuntimeError(
                f"Could not fetch YouTube transcript: {str(e)}. "
                "Make sure the video has captions/subtitles enabled."
            )

    def process_file(self, file_bytes: bytes, file_name: str) -> str:
        """Auto-detect file type and extract text."""
        name_lower = file_name.lower()
        if name_lower.endswith(".pdf"):
            return self.process_pdf(file_bytes)
        elif name_lower.endswith(".docx"):
            return self.process_docx(file_bytes)
        elif name_lower.endswith(".pptx"):
            return self.process_pptx(file_bytes)
        elif name_lower.endswith((".txt", ".md", ".csv")):
            return self.process_txt(file_bytes)
        else:
            # Try as plain text fallback
            return self.process_txt(file_bytes)

    def get_chunks(self, text: str) -> list[str]:
        """Split extracted text into chunks for embedding."""
        if not text.strip():
            return []
        return _chunk_text(text)